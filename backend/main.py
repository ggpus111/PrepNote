from dotenv import load_dotenv
load_dotenv()

import os
import json
import io
import re
from typing import Optional, Literal, List

from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

import pdfplumber
from docx import Document
from openai import OpenAI

# ✅ PPTX 지원
from pptx import Presentation as PptxPresentation

# ✅ OCR utils
from ocr_utils import ocr_image_bytes, ocr_images, pdf_to_images_pymupdf

app = FastAPI(title="PrepNote API", version="0.1.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "http://127.0.0.1:3000",
        "http://localhost:5173",
        "http://127.0.0.1:5173",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

client = OpenAI()
MODEL_TEXT = os.getenv("OPENAI_MODEL", "gpt-4.1-mini")


def _require_openai_key():
    if not os.getenv("OPENAI_API_KEY"):
        raise HTTPException(status_code=500, detail="OPENAI_API_KEY가 설정되지 않았습니다. (환경변수 확인)")


@app.get("/health")
def health():
    return {"ok": True}


# -----------------------------
# 파일 텍스트 추출 + OCR fallback
# -----------------------------
class ExtractTextResponse(BaseModel):
    text: str


def _clean_text(s: str) -> str:
    s = (s or "").strip()
    s = re.sub(r"[ \t]+\n", "\n", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def _is_text_enough(text: str, min_chars: int = 80) -> bool:
    return len((text or "").strip()) >= min_chars


def _extract_pdf_text(file_bytes: bytes) -> str:
    texts: List[str] = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages[:10]:
            t = page.extract_text() or ""
            if t.strip():
                texts.append(t.strip())
    return _clean_text("\n\n".join(texts))


def _extract_docx_text(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    lines = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    return _clean_text("\n".join(lines))


def _extract_pptx_text(file_bytes: bytes) -> str:
    prs = PptxPresentation(io.BytesIO(file_bytes))
    texts: List[str] = []
    for slide in prs.slides:
        slide_lines: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    slide_lines.append(t)
        if slide_lines:
            texts.append("\n".join(slide_lines))
    return _clean_text("\n\n".join(texts))


def _ocr_pdf_fallback(file_bytes: bytes) -> str:
    # 스캔 PDF 등: PDF -> 이미지 렌더 -> OCR
    images = pdf_to_images_pymupdf(file_bytes, max_pages=8, zoom=2.0)
    text = ocr_images(images, lang="eng+kor")
    return _clean_text(text)


def _ocr_pptx_images_fallback(file_bytes: bytes) -> str:
    # 이미지 위주 PPT: 슬라이드 내 picture blob OCR
    prs = PptxPresentation(io.BytesIO(file_bytes))
    img_blobs: List[bytes] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            # picture type (python-pptx: shape_type == 13)
            if getattr(shape, "shape_type", None) == 13:
                try:
                    img_blobs.append(shape.image.blob)
                except Exception:
                    pass
        if len(img_blobs) >= 8:
            break

    chunks: List[str] = []
    for b in img_blobs[:8]:
        try:
            t = ocr_image_bytes(b, lang="eng+kor")
            if t.strip():
                chunks.append(t.strip())
        except Exception:
            continue

    return _clean_text("\n\n".join(chunks))


@app.post("/extract-text", response_model=ExtractTextResponse)
async def extract_text(file: UploadFile = File(...)):
    filename = (file.filename or "").strip()
    content_type = (file.content_type or "").strip().lower()
    data = await file.read()

    if not filename:
        raise HTTPException(status_code=400, detail="파일 이름이 없습니다.")
    if not data:
        raise HTTPException(status_code=400, detail="업로드된 파일이 비어 있습니다.")

    name = filename.lower()

    # TXT
    if name.endswith(".txt") or content_type == "text/plain":
        try:
            text = data.decode("utf-8", errors="ignore").strip()
        except Exception:
            text = data.decode(errors="ignore").strip()

        text = _clean_text(text)
        if not text:
            raise HTTPException(status_code=400, detail="TXT에서 텍스트를 찾지 못했습니다.")
        return ExtractTextResponse(text=text)

    # 이미지 파일 (추가)
    if name.endswith((".png", ".jpg", ".jpeg", ".webp")):
        try:
            text = _clean_text(ocr_image_bytes(data, lang="eng+kor"))
            if not text:
                raise HTTPException(status_code=400, detail="이미지에서 텍스트를 추출하지 못했습니다.")
            return ExtractTextResponse(text=text)
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"이미지 OCR 실패: {str(e)}")

    # PDF
    if name.endswith(".pdf") or content_type == "application/pdf":
        try:
            text = _extract_pdf_text(data)

            # ✅ 텍스트가 부족하면 OCR fallback
            if not _is_text_enough(text):
                ocr_text = _ocr_pdf_fallback(data)
                # OCR 결과가 있으면 우선 사용, 아니면 기존(빈/짧은) 텍스트 유지
                if _is_text_enough(ocr_text, 30):
                    text = ocr_text

            if not text:
                raise HTTPException(
                    status_code=400,
                    detail="PDF에서 텍스트를 추출하지 못했습니다. (이미지/스캔 PDF일 수 있음)"
                )
            return ExtractTextResponse(text=text)
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"PDF 텍스트 추출 실패: {str(e)}")

    # DOCX
    if name.endswith(".docx") or content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        try:
            text = _extract_docx_text(data)
            if not text:
                raise HTTPException(status_code=400, detail="DOCX에서 텍스트를 추출하지 못했습니다.")
            return ExtractTextResponse(text=text)
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"DOCX 텍스트 추출 실패: {str(e)}")

    # PPTX
    if name.endswith(".pptx") or content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        try:
            text = _extract_pptx_text(data)

            # ✅ 텍스트가 부족하면 이미지 OCR fallback
            if not _is_text_enough(text):
                ocr_text = _ocr_pptx_images_fallback(data)
                if _is_text_enough(ocr_text, 30):
                    # 기존 텍스트가 조금이라도 있으면 합치기
                    text = _clean_text((text + "\n\n" + ocr_text).strip())

            if not text:
                raise HTTPException(
                    status_code=400,
                    detail="PPTX에서 텍스트를 추출하지 못했습니다. (이미지 위주 슬라이드일 수 있음)"
                )
            return ExtractTextResponse(text=text)
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"PPTX 텍스트 추출 실패: {str(e)}")

    raise HTTPException(
        status_code=400,
        detail=f"지원 형식: TXT / PDF / DOCX / PPTX (+ PNG/JPG/WEBP) (filename={filename}, content_type={content_type})"
    )


# -----------------------------
# 요약 생성 (표 + 난이도 질문 + 직장인 one-pager)
# -----------------------------
class SummaryOptions(BaseModel):
    length: Optional[Literal["short", "medium", "long"]] = "medium"
    audience: Optional[Literal["elementary", "middle", "high", "college", "office"]] = "college"


class SummaryRequest(BaseModel):
    title: Optional[str] = None
    text: str
    options: Optional[SummaryOptions] = None


class SummaryResponse(BaseModel):
    summaryId: str
    summary: str
    outline: List[str] = []


def _safe_json(text: str) -> Optional[dict]:
    if not text:
        return None
    text = text.strip()

    if "{" in text and "}" in text:
        start = text.find("{")
        end = text.rfind("}")
        if start != -1 and end != -1 and end > start:
            candidate = text[start:end + 1].strip()
            try:
                return json.loads(candidate)
            except Exception:
                return None
    return None


def _audience_profile(aud: str) -> str:
    return {
        "elementary": (
            "초등학생: 아주 쉬운 말, 짧은 문장. 어려운 용어 금지.\n"
            "영어(뜻) 형식으로 괄호 해석 필수. 예: apple(사과)\n"
            "예시는 1~2개만."
        ),
        "middle": (
            "중학생: 쉬운 말 + 핵심 개념 중심.\n"
            "영어(뜻) 괄호 해석을 자주 붙이되 과도하게 길게 설명하지 말 것."
        ),
        "high": (
            "고등학생: 개념+원리까지. 필요한 용어는 짧게 정의.\n"
            "영어는 등장 시 1회 정도만 (뜻) 병기."
        ),
        "college": (
            "대학생: 구조적/논리적 정리. 핵심 주장-근거-결론 흐름."
        ),
        "office": (
            "직장인: 결론/의사결정 중심. 핵심 요지 먼저.\n"
            "실무 적용, 리스크/한계, 액션 아이템을 명확히.\n"
            "회의/보고서 형태(불릿, 짧은 문장)."
        ),
    }.get(aud, "대학생: 구조적/논리적 정리.")


def _length_profile(length: str) -> str:
    return {
        "short": "짧게: 핵심만.",
        "medium": "보통: 공부/발표에 바로 쓸 수 있을 정도.",
        "long": "길게: 이해를 돕는 설명까지 포함.",
    }.get(length, "보통: 충분히.")


def _extract_internal_sources(text: str) -> List[str]:
    if not text:
        return []

    sources = set()

    for m in re.findall(r"(https?://[^\s)]+)", text):
        if len(m) > 10:
            sources.add(m.strip())

    patterns = [
        r"참고문헌[:：].*",
        r"참고 자료[:：].*",
        r"출처[:：].*",
        r"References[:：].*",
        r"Bibliography[:：].*",
        r"Source[:：].*",
    ]
    for p in patterns:
        for m in re.findall(p, text, flags=re.IGNORECASE):
            line = m.strip()
            if 5 < len(line) < 300:
                sources.add(line)

    return list(sources)[:12]


def _build_search_queries(title: str, text: str, aud: str) -> List[str]:
    head = (text or "")[:1200]
    tokens = re.findall(r"[A-Za-z]{4,}|[가-힣]{2,}", head)
    freq = {}
    for t in tokens:
        t = t.lower()
        if t in ["그리고", "하지만", "따라서", "또한", "이것", "저것", "대한", "관련", "있다", "한다"]:
            continue
        freq[t] = freq.get(t, 0) + 1

    top = sorted(freq.items(), key=lambda x: x[1], reverse=True)[:8]
    keywords = [k for k, _ in top]

    base = title.strip() if title else "주제"
    queries = []

    if keywords:
        queries.append(f"{base} {keywords[0]} 개념 정리")
        if len(keywords) > 1:
            queries.append(f"{base} {keywords[0]} {keywords[1]} 비교")
        if len(keywords) > 2:
            queries.append(f"{base} {keywords[2]} 정의")
        queries.append(f"{base} 발표 자료 요약")
        queries.append(f"{base} 참고문헌")
    else:
        queries = [f"{base} 개념 정리", f"{base} 발표 자료", f"{base} 참고문헌"]

    if aud in ["elementary", "middle"]:
        queries.insert(0, f"{base} 쉬운 설명")

    return queries[:6]


@app.post("/summaries", response_model=SummaryResponse)
def create_summary(req: SummaryRequest):
    _require_openai_key()

    title = (req.title or "제목 없음").strip()
    length = (req.options.length if req.options else "medium") or "medium"
    audience = (req.options.audience if req.options else "college") or "college"

    aud_rule = _audience_profile(audience)
    len_rule = _length_profile(length)

    internal_sources = _extract_internal_sources(req.text)
    queries = _build_search_queries(title, req.text, audience)

    system_prompt = f"""
너는 '공부/발표/회의에 바로 쓰는 자료 정리' 전문가다.
반드시 아래 규칙을 지켜 출력한다.

[대상/난이도]
{aud_rule}

[공통 규칙]
- 한국어로 작성
- 자료에 없는 내용은 만들지 말 것(추측 금지)
- 문장을 짧고 명확하게
- "대본"처럼 말하지 말고, '정리본' 문 reveal
- 초/중학생 대상이면 영어 단어는 apple(사과) 형태로 병기
- 표는 마크다운 표로 작성
- 출력은 반드시 JSON만(그 외 텍스트 금지)

[출처 처리]
- 문서 내부에 URL/참고문헌/출처 표기가 있으면 마지막에 '문서 내 출처'로 모아서 정리
- 문서 내부 출처가 없다면 '출처 찾기용 검색어' 제안(링크 지어내지 말 것)

[중요]
- 직장인(office) 대상일 때만 One-pager 섹션을 추가한다.
""".strip()

    one_pager_rule = ""
    if audience == "office":
        one_pager_rule = """
6) 회의용 1페이지 요약(One-pager)
- 목적
- 핵심 결론(1~3줄)
- 주요 근거(3~5개 bullet)
- 리스크/주의사항(2~4개 bullet)
- 다음 액션 아이템(To-do)(3~5개, 주어+동사)
""".strip()

    user_prompt = f"""
[제목]
{title}

[분량]
{len_rule}

[원문]
{req.text[:22000]}

[문서 내에서 발견한 출처 후보(있을 때만 참고)]
{json.dumps(internal_sources, ensure_ascii=False)}

[출처 찾기용 검색어 후보(출처가 없을 때 활용)]
{json.dumps(queries, ensure_ascii=False)}

[반드시 지켜야 할 출력 JSON 형식]
{{
  "summary": "아래 섹션 구조를 그대로 포함한 정리본 텍스트(줄바꿈 포함 가능)",
  "outline": ["슬라이드 제목 1", "슬라이드 제목 2", "..."]
}}

[summary 섹션 구조(이 순서/제목 그대로)]
1) 한눈에 요약(핵심 5~8줄)

2) 내용 정리 표(마크다운 표)
| 개념 | 정의 | 예시 | 주의점 |
|---|---|---|---|
- 핵심 개념 4~7개

3) 핵심 용어/키워드(정의 1줄씩, 5~10개)

4) 예상 질문(난이도별)
- 쉬움 2~3개
- 보통 2~3개
- 어려움 2~3개
(각각 질문 + 한 줄 답)

5) (있으면) 문서 내 출처 모음
   (없으면) 출처 찾기용 검색어(3~6개)

{one_pager_rule}

[outline 규칙]
- 7~10개 권장
- 서론-본론-결론 흐름
- 표(내용 정리)와 대응되게 중간 파트 구성
""".strip()

    r = client.responses.create(
        model=MODEL_TEXT,
        input=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
    )

    out = (r.output_text or "").strip()
    data = _safe_json(out)
    if not data:
        return SummaryResponse(
            summaryId="summary_" + os.urandom(4).hex(),
            summary=out,
            outline=[],
        )

    summary_text = str(data.get("summary", "")).strip()
    outline = data.get("outline", [])
    if not isinstance(outline, list):
        outline = []
    outline = [str(x).strip() for x in outline if str(x).strip()][:12]

    return SummaryResponse(
        summaryId="summary_" + os.urandom(4).hex(),
        summary=summary_text,
        outline=outline,
    )


# -----------------------------
# 대본 생성 (품질 업그레이드)
# -----------------------------
class ScriptOptions(BaseModel):
    tone: Optional[Literal["formal", "friendly", "energetic"]] = "formal"
    speakerCount: Optional[int] = 1
    targetMinutes: Optional[int] = 3


class ScriptRequest(BaseModel):
    title: str
    summaryText: str
    outline: Optional[List[str]] = None
    options: Optional[ScriptOptions] = None


class ScriptResponse(BaseModel):
    scriptId: str
    content: List[str]


@app.post("/scripts", response_model=ScriptResponse)
def create_script(req: ScriptRequest):
    _require_openai_key()

    opt = req.options or ScriptOptions()
    speaker_count = max(1, min(int(opt.speakerCount or 1), 5))
    minutes = max(1, min(int(opt.targetMinutes or 3), 15))
    tone = (opt.tone or "formal").strip()

    tone_desc = {
        "formal": "존댓말/발표체/깔끔하게",
        "friendly": "부드럽고 친근한 존댓말",
        "energetic": "또렷하고 자신감 있게(과장 금지)",
    }.get(tone, "존댓말/발표체/깔끔하게")

    outline_text = "\n".join([f"- {x}" for x in (req.outline or [])]) if req.outline else ""

    system_prompt = f"""
너는 '실제로 발표 가능한 대본'을 만드는 조교다.
반드시 아래 규칙을 지켜라.

[공통]
- 한국어
- 톤: {tone_desc}
- 목표 시간: {minutes}분 (시간에 맞게 분량 조절)
- 발표자 {speaker_count}명 → 발표자별로 자연스럽게 역할 분담
- 각 대본은 반드시 "발표자 N:"으로 시작
- 발표 흐름: 도입(관심 끌기) → 목차 안내 → 핵심 내용(정리) → 요약/결론 → 마무리 멘트
- 슬라이드/목차가 있으면 슬라이드 단위로 전환 멘트 포함(예: "다음은 ~입니다")
- 과장/환각 금지: 요약 텍스트 밖의 사실을 만들어내지 말 것
- 출력은 JSON만
""".strip()

    user_prompt = f"""
[발표 제목]
{req.title}

[권장 목차/슬라이드]
{outline_text if outline_text else "(없음: 요약 텍스트 구조를 기준으로 자연스럽게 목차를 만들어도 됨)"}

[요약/원문 기반 텍스트]
{req.summaryText[:20000]}

[출력 JSON 형식]
{{
  "content": [
    "발표자 1: ...",
    "발표자 2: ..."
  ]
}}

[추가 요구]
- {minutes}분에 맞게 핵심만 남기되, 뜬금없는 생략 없이 '설명 → 예시 한 번 → 정리' 리듬 유지
- 발표자가 여러 명이면: 중복 설명 금지, 각자 구간이 명확해야 함
""".strip()

    r = client.responses.create(
        model=MODEL_TEXT,
        input=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
    )
    out = (r.output_text or "").strip()

    data = _safe_json(out)
    if not data:
        return ScriptResponse(
            scriptId="script_" + os.urandom(4).hex(),
            content=[out] if out else ["(대본 생성 결과가 비어있습니다)"],
        )

    content = data.get("content", [])
    if not isinstance(content, list):
        content = []

    content = [str(x).strip() for x in content if str(x).strip()]
    if not content:
        content = ["(대본 생성 결과가 비어있습니다)"]

    return ScriptResponse(
        scriptId="script_" + os.urandom(4).hex(),
        content=content,
    )
